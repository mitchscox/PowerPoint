# Fixing the issue by using `add_paragraph().text` instead of passing argument directly.

from pptx import Presentation

# Create a presentation object
prs = Presentation()

# Slide 1 - Title Slide
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]
title.text = "Architecting a Software Testing Solution"
subtitle.text = "For Multiple Financial Systems"

# Slide 2 - Introduction
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
title_2 = slide_2.shapes.title
title_2.text = "Overview"
bullet_points_2 = slide_2.shapes.placeholders[1].text_frame
bullet_points_2.text = "Steps to architect a testing solution:"
bullet_points_2.add_paragraph().text = "1. Understand the system landscape"
bullet_points_2.add_paragraph().text = "2. Define testing scope"
bullet_points_2.add_paragraph().text = "3. Develop test strategy"
bullet_points_2.add_paragraph().text = "4. Select tools"
bullet_points_2.add_paragraph().text = "5. Ensure compliance & security"
bullet_points_2.add_paragraph().text = "6. Set up test environments"
bullet_points_2.add_paragraph().text = "7. Collaborate effectively"

# Slide 3 - Understand the System Landscape
slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
title_3 = slide_3.shapes.title
title_3.text = "Understand the System Landscape"
bullet_points_3 = slide_3.shapes.placeholders[1].text_frame
bullet_points_3.text = "Key aspects:"
bullet_points_3.add_paragraph().text = "• Inventory of systems"
bullet_points_3.add_paragraph().text = "• System interactions and dependencies"
bullet_points_3.add_paragraph().text = "• Technology stack compatibility"

# Slide 4 - Define Testing Scope
slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
title_4 = slide_4.shapes.title
title_4.text = "Define Testing Scope"
bullet_points_4 = slide_4.shapes.placeholders[1].text_frame
bullet_points_4.text = "Types of testing:"
bullet_points_4.add_paragraph().text = "• Functional testing"
bullet_points_4.add_paragraph().text = "• Integration testing"
bullet_points_4.add_paragraph().text = "• Non-functional testing (performance, security)"
bullet_points_4.add_paragraph().text = "• Compliance testing"

# Slide 5 - Develop Test Strategy
slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
title_5 = slide_5.shapes.title
title_5.text = "Test Strategy & Coverage"
bullet_points_5 = slide_5.shapes.placeholders[1].text_frame
bullet_points_5.text = "Focus areas:"
bullet_points_5.add_paragraph().text = "• Risk-based testing"
bullet_points_5.add_paragraph().text = "• End-to-end workflows"
bullet_points_5.add_paragraph().text = "• Real data considerations"

# Slide 6 - Tool Selection
slide_6 = prs.slides.add_slide(prs.slide_layouts[1])
title_6 = slide_6.shapes.title
title_6.text = "Tool Selection"
bullet_points_6 = slide_6.shapes.placeholders[1].text_frame
bullet_points_6.text = "Testing and automation tools:"
bullet_points_6.add_paragraph().text = "• Automation: Selenium, RestAssured"
bullet_points_6.add_paragraph().text = "• Test management: Jira, TestRail"
bullet_points_6.add_paragraph().text = "• CI integration: Jenkins, GitLab CI"
bullet_points_6.add_paragraph().text = "• Monitoring: Real-time dashboards"

# Slide 7 - Compliance & Security
slide_7 = prs.slides.add_slide(prs.slide_layouts[1])
title_7 = slide_7.shapes.title
title_7.text = "Compliance & Security"
bullet_points_7 = slide_7.shapes.placeholders[1].text_frame
bullet_points_7.text = "Key considerations:"
bullet_points_7.add_paragraph().text = "• Industry standards: PCI-DSS, SOC 2"
bullet_points_7.add_paragraph().text = "• Penetration testing & vulnerability assessments"

# Slide 8 - Environment Setup
slide_8 = prs.slides.add_slide(prs.slide_layouts[1])
title_8 = slide_8.shapes.title
title_8.text = "Environment Setup"
bullet_points_8 = slide_8.shapes.placeholders[1].text_frame
bullet_points_8.text = "Test environment setup:"
bullet_points_8.add_paragraph().text = "• Sandbox/staging environments"
bullet_points_8.add_paragraph().text = "• Docker/Kubernetes for isolated, scalable environments"

# Slide 9 - Collaboration
slide_9 = prs.slides.add_slide(prs.slide_layouts[1])
title_9 = slide_9.shapes.title
title_9.text = "Collaboration"
bullet_points_9 = slide_9.shapes.placeholders[1].text_frame
bullet_points_9.text = "Team collaboration:"
bullet_points_9.add_paragraph().text = "• DevOps integration for CI/CD"
bullet_points_9.add_paragraph().text = "• Business analysts for real-world scenarios"

# Save the presentation
pptx_path = "/home/bugeye2/Testing_Solution_Architecture.pptx"
prs.save(pptx_path)

pptx_path
